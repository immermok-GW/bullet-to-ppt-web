import streamlit as st

st.set_page_config(
    page_title="Bullet to PPT 神器",          # 瀏覽器標籤名稱
    page_icon="🚀",                           # 頁面小圖示（可以用 emoji）
    layout="wide",                            # 寬版布局（左右留白變少，更專業）
    initial_sidebar_state="auto"              # 側邊欄自動展開/收合
)

import json
from openai import OpenAI
from pptx import Presentation
from pptx.util import Pt
import io
import uuid

st.set_page_config(page_title="Bullet → PPT 神器", layout="centered")
st.title("🚀 Bullet Points 一鍵生成專業 PPT")
st.markdown("輸入主題 + 每行一個 bullet point，AI 自動做出 8 頁簡報！")

theme = st.text_input("📌 PPT 主題", "2026 香港樓市分析")
bullets = st.text_area("📝 Bullet Points（每行一個）", 
                       "香港樓價趨勢\n政府新政策\n投資建議\n風險分析", 
                       height=200)

if st.button("✨ 生成 PPT", type="primary", use_container_width=True):
    if not theme or not bullets:
        st.error("請填寫主題和 bullet points")
    else:
        with st.spinner("AI 正在生成 PPT...（約 10–20 秒）"):
            try:
                client = OpenAI(
                    api_key=st.secrets["DEEPSEEK_KEY"],
                    base_url="https://api.deepseek.com"
                )
                
                prompt = f"""你是專業 PPT 設計師。主題：{theme}
                Bullet Points：
                {bullets}
                請輸出純 JSON 格式（8 頁），每頁包含：
                {{"title": "標題", "content": ["bullet1", "bullet2", ...]}}
                直接回傳 JSON，不要多餘文字。"""
                
                response = client.chat.completions.create(
                    model="deepseek-chat",
                    messages=[{"role": "user", "content": prompt}]
                )
                
                slides = json.loads(response.choices[0].message.content.strip())
                
                prs = Presentation()
                for slide in slides:
                    s = prs.slides.add_slide(prs.slide_layouts[1])
                    s.shapes.title.text = slide["title"]
                    tf = s.shapes.placeholders[1].text_frame
                    tf.clear()
                    for txt in slide.get("content", []):
                        p = tf.add_paragraph()
                        p.text = txt
                        p.font.size = Pt(20)
                
                bio = io.BytesIO()
                prs.save(bio)
                bio.seek(0)
                
                st.success("✅ 生成成功！")
                st.download_button(
                    label="📥 下載你的 .pptx 檔案",
                    data=bio,
                    file_name=f"{theme}_{uuid.uuid4().hex[:8]}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )
                st.balloons()
            except Exception as e:
                st.error(f"出錯了：{str(e)}")
